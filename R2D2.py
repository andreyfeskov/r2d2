import win32com.client
from pptx import Presentation
from pptx.util import Inches

source_path = '1.pptx'
word_document_path = 'C:\\DEN2\\me.docx'


prs = Presentation(source_path)

word = win32com.client.Dispatch("Word.Application")
word.Documents.Open(word_document_path)

text = word.ActiveDocument.Content.Text
paragraphs = text.split('\n\n')

slide_count = len(prs.slides)
current_slide = 1
for paragraph in paragraphs:
    if current_slide > slide_count:
        # Создаем новый слайд, если достигли конца существующих слайдов
        prs.slides.add_slide(prs.slide_layouts[0])
        current_slide += 1
    try:
        title_placeholder = prs.slides[current_slide-1].shapes.title
        body_placeholder = prs.slides[current_slide-1].placeholders[1]

        title_placeholder.text = paragraph.split('\n')[0]
        body_placeholder.text = '\n'.join(paragraph.split('\n')[1:])
    except IndexError:
        pass
    current_slide += 1
prs.save('NEW_презентация.pptx')

word.ActiveDocument.Close()
word.Quit()
